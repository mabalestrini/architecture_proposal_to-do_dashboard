#!/usr/bin/env python3
"""Extract text from a PowerPoint (.pptx) file."""

import sys
from pathlib import Path
from typing import Optional

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx is not installed. Run: pip install python-pptx")
    sys.exit(1)


def extract_text(pptx_path: str, output_path: Optional[str] = None) -> None:
    path = Path(pptx_path)
    if not path.exists():
        print(f"Error: file not found: {pptx_path}")
        sys.exit(1)

    prs = Presentation(path)
    lines = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {slide_num} ---")

        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    row_cells = []
                    for cell in row.cells:
                        cell_text = cell.text_frame.text.strip()
                        if cell_text:
                            row_cells.append(cell_text)
                    if row_cells:
                        lines.append("\t".join(row_cells))
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("[Notes]")
                lines.append(notes_text)

        lines.append("")

    output = "\n".join(lines)

    if output_path:
        Path(output_path).write_text(output, encoding="utf-8")
        print(f"Text written to {output_path}")
    else:
        print(output)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(f"Usage: {sys.argv[0]} <file.pptx> [output.txt]")
        sys.exit(1)

    extract_text(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
